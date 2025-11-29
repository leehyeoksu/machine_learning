#!/usr/bin/env python3
"""
커리어 로드맵 2025-2029: 성장 전략 프레젠테이션
Career Roadmap Presentation Generator
"""

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
except ImportError:
    print("python-pptx 라이브러리를 먼저 설치해주세요:")
    print("pip install python-pptx")
    exit(1)

# 프레젠테이션 생성
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# 슬라이드 1: 타이틀
slide1 = prs.slides.add_slide(prs.slide_layouts[6])
background = slide1.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(20, 30, 48)

title_box = slide1.shapes.add_textbox(Inches(1), Inches(2.3), Inches(8), Inches(1.5))
title_frame = title_box.text_frame
title_frame.text = "커리어 로드맵"
title_paragraph = title_frame.paragraphs[0]
title_paragraph.font.size = Pt(60)
title_paragraph.font.bold = True
title_paragraph.font.color.rgb = RGBColor(255, 255, 255)
title_paragraph.alignment = PP_ALIGN.CENTER

subtitle_box = slide1.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(0.8))
subtitle_frame = subtitle_box.text_frame
subtitle_frame.text = "2025-2029 성장 전략"
subtitle_paragraph = subtitle_frame.paragraphs[0]
subtitle_paragraph.font.size = Pt(32)
subtitle_paragraph.font.color.rgb = RGBColor(180, 200, 230)
subtitle_paragraph.alignment = PP_ALIGN.CENTER

year_box = slide1.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(0.6))
year_frame = year_box.text_frame
year_frame.text = "부트캠프 → 학부연구생 → 대학원"
year_paragraph = year_frame.paragraphs[0]
year_paragraph.font.size = Pt(24)
year_paragraph.font.color.rgb = RGBColor(255, 215, 0)
year_paragraph.alignment = PP_ALIGN.CENTER

# 슬라이드 2: 전체 로드맵 개요
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
background = slide2.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(245, 248, 250)

title_box = slide2.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
title_frame = title_box.text_frame
title_frame.text = "전체 로드맵 개요 (2025-2029)"
title_paragraph = title_frame.paragraphs[0]
title_paragraph.font.size = Pt(36)
title_paragraph.font.bold = True
title_paragraph.font.color.rgb = RGBColor(20, 30, 48)

content_box = slide2.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(8.4), Inches(5.3))
text_frame = content_box.text_frame
text_frame.word_wrap = True

phases = [
    ("Phase 1: 부트캠프 (2025.11 ~ 2026.04)", "이스트소프트 AI 부트캠프\nML/DL 기초 + PyTorch + FastAPI 심화", RGBColor(66, 133, 244)),
    ("Phase 2: 학부연구생 (2026.05 ~ 2026.12)", "경희대 연구실 합류\n연구 경험 + 포트폴리오 구축", RGBColor(52, 168, 83)),
    ("Phase 3: 대학원 진학 (2027.02 입학)", "3가지 전략적 옵션\n자대 / 타대(SPK) / 타대(YK/IST)", RGBColor(251, 188, 5)),
    ("Phase 4: 석사과정 (2027.03 ~ 2029.02)", "연구 성과 창출 + 논문 발표\n취업 준비 병행", RGBColor(234, 67, 53))
]

for i, (phase, desc, color) in enumerate(phases):
    p = text_frame.add_paragraph() if i > 0 else text_frame.paragraphs[0]
    p.text = phase
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = color
    p.space_after = Pt(6)
    if i > 0:
        p.space_before = Pt(18)
    
    for line in desc.split('\n'):
        p = text_frame.add_paragraph()
        p.text = f"  • {line}"
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(60, 60, 60)
        p.level = 1
        p.space_after = Pt(4)

# 슬라이드 3: Phase 1 - 부트캠프
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
background = slide3.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(232, 240, 254)

title_box = slide3.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
title_frame = title_box.text_frame
title_frame.text = "Phase 1: 부트캠프 (2025.11 ~ 2026.04)"
title_paragraph = title_frame.paragraphs[0]
title_paragraph.font.size = Pt(36)
title_paragraph.font.bold = True
title_paragraph.font.color.rgb = RGBColor(25, 103, 210)

subtitle_box = slide3.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(0.6))
subtitle_frame = subtitle_box.text_frame
subtitle_frame.text = "🎯 목표: 이론과 실전 역량 균형 강화"
subtitle_paragraph = subtitle_frame.paragraphs[0]
subtitle_paragraph.font.size = Pt(24)
subtitle_paragraph.font.bold = True
subtitle_paragraph.font.color.rgb = RGBColor(66, 133, 244)

content_box = slide3.shapes.add_textbox(Inches(0.8), Inches(2.3), Inches(8.4), Inches(4.7))
text_frame = content_box.text_frame
text_frame.word_wrap = True

sections = [
    ("핵심 학습 영역", [
        "ML/DL 기초 이론 정리 (SVM, 경사하강법, Convex Optimization)",
        "PyTorch 딥러닝 실습 심화 (논문 구현, 오픈소스 따라하기)",
        "FastAPI + Flask 웹 프레임워크 마스터",
        "Docker + Celery를 활용한 서버 사이드 프로그래밍"
    ]),
    ("필수 프로젝트", [
        "End-to-End ML 프로젝트 (데이터 준비 → 학습 → 배포)",
        "PyTorch + FastAPI 통합 서비스 구축",
        "MLOps 도구 경험 (Docker, GitHub Actions)",
        "부트캠프 팀 프로젝트 → 포트폴리오화"
    ]),
    ("학습 우선순위", [
        "1️⃣ 수학/이론 보완 (선형대수, 확률론)",
        "2️⃣ 기술스택 심화 (PyTorch, Docker, Redis)",
        "3️⃣ 통합 프로젝트 경험 (주/월 단위 결과물 생산)"
    ])
]

for i, (section_title, points) in enumerate(sections):
    p = text_frame.add_paragraph() if i > 0 else text_frame.paragraphs[0]
    p.text = f"📌 {section_title}"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(25, 103, 210)
    p.space_after = Pt(6)
    if i > 0:
        p.space_before = Pt(14)
    
    for point in points:
        p = text_frame.add_paragraph()
        p.text = f"• {point}"
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(40, 40, 40)
        p.level = 1
        p.space_after = Pt(3)

# 슬라이드 4: Phase 2 - 학부연구생
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
background = slide4.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(232, 245, 233)

title_box = slide4.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
title_frame = title_box.text_frame
title_frame.text = "Phase 2: 학부연구생 (2026.05 ~ 2026.12)"
title_paragraph = title_frame.paragraphs[0]
title_paragraph.font.size = Pt(36)
title_paragraph.font.bold = True
title_paragraph.font.color.rgb = RGBColor(27, 94, 32)

subtitle_box = slide4.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(0.6))
subtitle_frame = subtitle_box.text_frame
subtitle_frame.text = "🎯 목표: 연구 수행 능력 & 팀 협업 역량 강화"
subtitle_paragraph = subtitle_frame.paragraphs[0]
subtitle_paragraph.font.size = Pt(24)
subtitle_paragraph.font.bold = True
subtitle_paragraph.font.color.rgb = RGBColor(52, 168, 83)

content_box = slide4.shapes.add_textbox(Inches(0.8), Inches(2.3), Inches(8.4), Inches(4.7))
text_frame = content_box.text_frame
text_frame.word_wrap = True

sections = [
    ("핵심 경험 요소", [
        "문제 정의 & 문헌 조사 (관련 논문 읽기 및 요약)",
        "연구 기획 & 실험 설계 (소규모 연구과제 주도)",
        "코딩 & 시스템 통합 (FastAPI, PyTorch, Docker 활용)",
        "결과 분석 & 보고 (연구실 세미나 발표)"
    ]),
    ("구체적 활동", [
        "최소 1개 연구 프로젝트 처음부터 끝까지 경험",
        "대학원생 멘토와 협업 (주간 미팅, 피드백 사이클)",
        "국내 학회 포스터 발표 또는 교내 논문경진대회 도전",
        "여름방학: AI 해커톤/공모전 참가 (수상 경력 노리기)"
    ]),
    ("대학원 진학 준비", [
        "영어 성적 준비 (TOEIC 900+ / TOEFL 90+)",
        "희망 대학원 리스트업 & 교수 컨택 시작",
        "추천서 확보 (연구실 교수님, 부트캠프 멘토)",
        "포트폴리오 프로젝트 3개+ 완성"
    ])
]

for i, (section_title, points) in enumerate(sections):
    p = text_frame.add_paragraph() if i > 0 else text_frame.paragraphs[0]
    p.text = f"📌 {section_title}"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(27, 94, 32)
    p.space_after = Pt(6)
    if i > 0:
        p.space_before = Pt(12)
    
    for point in points:
        p = text_frame.add_paragraph()
        p.text = f"• {point}"
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(40, 40, 40)
        p.level = 1
        p.space_after = Pt(3)

# 슬라이드 5: Phase 3 개요
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
background = slide5.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 248, 225)

title_box = slide5.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
title_frame = title_box.text_frame
title_frame.text = "Phase 3: 대학원 진학 전략 (2027.02 입학)"
title_paragraph = title_frame.paragraphs[0]
title_paragraph.font.size = Pt(36)
title_paragraph.font.bold = True
title_paragraph.font.color.rgb = RGBColor(245, 124, 0)

subtitle_box = slide5.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(0.6))
subtitle_frame = subtitle_box.text_frame
subtitle_frame.text = "🎯 3가지 전략적 옵션"
subtitle_paragraph = subtitle_frame.paragraphs[0]
subtitle_paragraph.font.size = Pt(28)
subtitle_paragraph.font.bold = True
subtitle_paragraph.font.color.rgb = RGBColor(230, 81, 0)

content_box = slide5.shapes.add_textbox(Inches(0.8), Inches(2.4), Inches(8.4), Inches(4.5))
text_frame = content_box.text_frame
text_frame.word_wrap = True

options = [
    ("옵션 1: 자대 (경희대) 진학", "[강력 추천 / 안전]", "학부연구생 활동 랩으로 진학\n장학금 & 인건비 협상 유리", RGBColor(46, 125, 50)),
    ("옵션 2: 타대 (SPK) 진학", "[도전 / 위험]", "4학년 여름방학 인턴십 프로그램 지원\nOpen Lab 통한 우회로 모색", RGBColor(211, 47, 47)),
    ("옵션 3: 타대 (YK/IST) 진학", "[도전 / 가능]", "GitHub + 기술블로그 + Cold Email\n공격적인 컨택 전략", RGBColor(245, 124, 0))
]

for i, (option, status, desc, color) in enumerate(options):
    p = text_frame.add_paragraph() if i > 0 else text_frame.paragraphs[0]
    p.text = option
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RGBColor(20, 30, 48)
    p.space_after = Pt(4)
    if i > 0:
        p.space_before = Pt(16)
    
    p = text_frame.add_paragraph()
    p.text = f"   {status}"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = color
    p.level = 1
    p.space_after = Pt(4)
    
    for line in desc.split('\n'):
        p = text_frame.add_paragraph()
        p.text = f"   • {line}"
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(60, 60, 60)
        p.level = 1
        p.space_after = Pt(3)

# 슬라이드 6: 옵션 1 - 자대
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
background = slide6.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(232, 245, 233)

title_box = slide6.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
title_frame = title_box.text_frame
title_frame.text = "옵션 1: 자대 (경희대) 진학"
title_paragraph = title_frame.paragraphs[0]
title_paragraph.font.size = Pt(36)
title_paragraph.font.bold = True
title_paragraph.font.color.rgb = RGBColor(27, 94, 32)

judgment_box = slide6.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(8.4), Inches(0.8))
text_frame = judgment_box.text_frame
p = text_frame.paragraphs[0]
p.text = "✅ [강력 추천 / 안전]"
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = RGBColor(46, 125, 50)

analysis_box = slide6.shapes.add_textbox(Inches(0.8), Inches(2.6), Inches(8.4), Inches(3.2))
text_frame = analysis_box.text_frame
text_frame.word_wrap = True

p = text_frame.paragraphs[0]
p.text = "전략적 판단"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = RGBColor(27, 94, 32)
p.space_after = Pt(10)

points = [
    "✓ 이미 검증된 인력으로 분류됨",
    "✓ 최상위 랩(Top Lab) 진학 가능성",
    "✓ 학부 연구생 경험 활용",
    "✓ 장학금 & 인건비 협상 유리",
    "✓ 지도교수와의 신뢰 관계 구축됨"
]

for point in points:
    p = text_frame.add_paragraph()
    p.text = point
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(40, 40, 40)
    p.space_after = Pt(8)
    p.level = 1

action_box = slide6.shapes.add_textbox(Inches(0.8), Inches(5.9), Inches(8.4), Inches(1.2))
text_frame = action_box.text_frame
text_frame.word_wrap = True
p = text_frame.paragraphs[0]
p.text = "🎯 Best Action"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = RGBColor(27, 94, 32)
p = text_frame.add_paragraph()
p.text = "학부 연구생 활동을 했던 랩으로 직행 진학"
p.font.size = Pt(20)
p.font.color.rgb = RGBColor(20, 30, 48)
p.level = 1

# 슬라이드 7: 옵션 2 - 타대 SPK
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
background = slide7.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 235, 238)

title_box = slide7.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
title_frame = title_box.text_frame
title_frame.text = "옵션 2: 타대 (SPK) 진학"
title_paragraph = title_frame.paragraphs[0]
title_paragraph.font.size = Pt(36)
title_paragraph.font.bold = True
title_paragraph.font.color.rgb = RGBColor(183, 28, 28)

judgment_box = slide7.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(8.4), Inches(0.8))
text_frame = judgment_box.text_frame
p = text_frame.paragraphs[0]
p.text = "⚠️ [도전 / 위험]"
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = RGBColor(211, 47, 47)

analysis_box = slide7.shapes.add_textbox(Inches(0.8), Inches(2.6), Inches(8.4), Inches(2.8))
text_frame = analysis_box.text_frame
text_frame.word_wrap = True

p = text_frame.paragraphs[0]
p.text = "전략적 판단"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = RGBColor(183, 28, 28)
p.space_after = Pt(10)

points = [
    "⚠ 학점 컷오프 존재 (높은 진입장벽)",
    "⚠ 획기적인 연구 실적 필수",
    "⚠ Top-tier 논문 없이는 합격 어려움",
    "⚠ 매우 높은 경쟁률"
]

for point in points:
    p = text_frame.add_paragraph()
    p.text = point
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(40, 40, 40)
    p.space_after = Pt(8)
    p.level = 1

action_box = slide7.shapes.add_textbox(Inches(0.8), Inches(5.5), Inches(8.4), Inches(1.6))
text_frame = action_box.text_frame
text_frame.word_wrap = True
p = text_frame.paragraphs[0]
p.text = "🎯 우회 전략"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = RGBColor(183, 28, 28)
p.space_after = Pt(8)

p = text_frame.add_paragraph()
p.text = "• 4학년 여름방학 (2026년 여름) 인턴십 프로그램 지원"
p.font.size = Pt(19)
p.font.color.rgb = RGBColor(20, 30, 48)
p.level = 1
p.space_after = Pt(4)

p = text_frame.add_paragraph()
p.text = "• Open Lab을 통한 우회로 모색"
p.font.size = Pt(19)
p.font.color.rgb = RGBColor(20, 30, 48)
p.level = 1

# 슬라이드 8: 옵션 3 - 타대 YK/IST 개요
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
background = slide8.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 243, 224)

title_box = slide8.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
title_frame = title_box.text_frame
title_frame.text = "옵션 3: 타대 (YK/IST) 진학"
title_paragraph = title_frame.paragraphs[0]
title_paragraph.font.size = Pt(36)
title_paragraph.font.bold = True
title_paragraph.font.color.rgb = RGBColor(230, 81, 0)

judgment_box = slide8.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(8.4), Inches(0.8))
text_frame = judgment_box.text_frame
p = text_frame.paragraphs[0]
p.text = "💪 [도전 / 가능]"
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = RGBColor(245, 124, 0)

insight_box = slide8.shapes.add_textbox(Inches(0.8), Inches(2.6), Inches(8.4), Inches(4.3))
text_frame = insight_box.text_frame
text_frame.word_wrap = True

p = text_frame.paragraphs[0]
p.text = "핵심 인사이트"
p.font.size = Pt(26)
p.font.bold = True
p.font.color.rgb = RGBColor(230, 81, 0)
p.space_after = Pt(14)

insights = [
    "시스템 분야 랩은 실무 능력 중시",
    "논문보다 구현 능력과 기술적 깊이가 중요",
    "효과적인 컨택이 합격의 핵심",
    "포트폴리오와 기술 역량 증명 필수"
]

for insight in insights:
    p = text_frame.add_paragraph()
    p.text = f"✓ {insight}"
    p.font.size = Pt(22)
    p.font.color.rgb = RGBColor(40, 40, 40)
    p.space_after = Pt(10)
    p.level = 1

p = text_frame.add_paragraph()
p.text = "🎯 핵심 3요소 전략"
p.font.size = Pt(26)
p.font.bold = True
p.font.color.rgb = RGBColor(230, 81, 0)
p.space_after = Pt(10)
p.space_before = Pt(16)

strategies = [
    "1️⃣ GitHub 포트폴리오 강화",
    "2️⃣ 기술 블로그 운영",
    "3️⃣ 공격적 Cold Email 전략"
]

for strategy in strategies:
    p = text_frame.add_paragraph()
    p.text = strategy
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(60, 60, 60)
    p.space_after = Pt(6)
    p.level = 1

# 슬라이드 9: GitHub 포트폴리오
slide9 = prs.slides.add_slide(prs.slide_layouts[6])
background = slide9.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(240, 248, 255)

title_box = slide9.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.7))
title_frame = title_box.text_frame
title_frame.text = "1️⃣ GitHub 포트폴리오 구축 전략"
title_paragraph = title_frame.paragraphs[0]
title_paragraph.font.size = Pt(32)
title_paragraph.font.bold = True
title_paragraph.font.color.rgb = RGBColor(13, 71, 161)

content_box = slide9.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(8.6), Inches(5.5))
text_frame = content_box.text_frame
text_frame.word_wrap = True

sections = [
    ("프로젝트 선정 기준", [
        "시스템 프로그래밍 (OS, Network, Distributed Systems)",
        "성능 최적화 사례 포함",
        "실제 문제 해결 중심 (Toy Project 지양)"
    ]),
    ("README 작성 필수 요소", [
        "Problem Statement (해결하려는 문제)",
        "Architecture Diagram (시스템 구조도)",
        "Performance Metrics (성능 지표)",
        "Challenges & Solutions (도전과 해결 과정)"
    ]),
    ("코드 품질 관리", [
        "일관된 코딩 스타일 (PEP8, Google Style Guide)",
        "의미 있는 주석 (Why, not What)",
        "Unit Test 포함 (테스트 커버리지 50%+)"
    ])
]

for i, (section_title, points) in enumerate(sections):
    p = text_frame.add_paragraph() if i > 0 else text_frame.paragraphs[0]
    p.text = f"📌 {section_title}"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RGBColor(13, 71, 161)
    p.space_after = Pt(8)
    if i > 0:
        p.space_before = Pt(16)
    
    for point in points:
        p = text_frame.add_paragraph()
        p.text = f"• {point}"
        p.font.size = Pt(17)
        p.font.color.rgb = RGBColor(40, 40, 40)
        p.level = 1
        p.space_after = Pt(4)

# 슬라이드 10: 기술 블로그
slide10 = prs.slides.add_slide(prs.slide_layouts[6])
background = slide10.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(245, 245, 250)

title_box = slide10.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.7))
title_frame = title_box.text_frame
title_frame.text = "2️⃣ 기술 블로그 운영 전략"
title_paragraph = title_frame.paragraphs[0]
title_paragraph.font.size = Pt(32)
title_paragraph.font.bold = True
title_paragraph.font.color.rgb = RGBColor(74, 20, 140)

content_box = slide10.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(8.6), Inches(5.5))
text_frame = content_box.text_frame
text_frame.word_wrap = True

sections = [
    ("추천 주제", [
        "시스템 성능 최적화 경험",
        "복잡한 버그 디버깅 과정",
        "운영체제 내부 동작 분석",
        "네트워크 프로토콜 구현 경험",
        "분산 시스템 설계 및 트레이드오프"
    ]),
    ("포스트 작성 구조", [
        "Problem: 무엇을 해결하려 했는가?",
        "Approach: 어떤 방법을 시도했는가?",
        "Deep Dive: 기술적으로 어떻게 구현했는가?",
        "Results: 결과와 성능 개선 수치",
        "Lessons Learned: 배운 점과 향후 개선"
    ]),
    ("전략적 팁", [
        "Medium / 개인 블로그 (dev.to, velog) 활용",
        "월 1-2회 정기 포스팅 (Consistency)",
        "영어 버전 병행 작성 (국제적 가시성)",
        "코드 스니펫 & 다이어그램 풍부하게"
    ])
]

for i, (section_title, points) in enumerate(sections):
    p = text_frame.add_paragraph() if i > 0 else text_frame.paragraphs[0]
    p.text = f"📝 {section_title}"
    p.font.size = Pt(21)
    p.font.bold = True
    p.font.color.rgb = RGBColor(74, 20, 140)
    p.space_after = Pt(8)
    if i > 0:
        p.space_before = Pt(14)
    
    for point in points:
        p = text_frame.add_paragraph()
        p.text = f"• {point}"
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(40, 40, 40)
        p.level = 1
        p.space_after = Pt(3)

# 슬라이드 11: Cold Email 전략
slide11 = prs.slides.add_slide(prs.slide_layouts[6])
background = slide11.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 248, 240)

title_box = slide11.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.7))
title_frame = title_box.text_frame
title_frame.text = "3️⃣ 공격적 Cold Email 전략"
title_paragraph = title_frame.paragraphs[0]
title_paragraph.font.size = Pt(32)
title_paragraph.font.bold = True
title_paragraph.font.color.rgb = RGBColor(191, 54, 12)

content_box = slide11.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(8.6), Inches(5.5))
text_frame = content_box.text_frame
text_frame.word_wrap = True

sections = [
    ("사전 준비 (Research)", [
        "교수님의 최근 3년 논문 리스트 확인",
        "연구실 홈페이지 진행 프로젝트 파악",
        "연구실 졸업생 진로 조사 (LinkedIn)",
        "최근 학회 발표 자료 확인"
    ]),
    ("이메일 구성 요소", [
        "Subject: 간결하고 구체적",
        "Opening: 특정 논문/프로젝트 언급",
        "Body: 관련 경험 & 역량 (GitHub 링크)",
        "Technical Blog: 깊이 있는 포스트 1-2개",
        "Closing: 간단한 미팅 요청 (15-30분)"
    ]),
    ("타이밍 & Follow-up", [
        "학기 시작 2-3개월 전 발송",
        "1주일 후 응답 없으면 1회 Follow-up",
        "여러 교수님께 동시 발송 (3-5명)",
        "거절 시 피드백 요청"
    ])
]

for i, (section_title, points) in enumerate(sections):
    p = text_frame.add_paragraph() if i > 0 else text_frame.paragraphs[0]
    p.text = f"✉️ {section_title}"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(191, 54, 12)
    p.space_after = Pt(7)
    if i > 0:
        p.space_before = Pt(12)
    
    for point in points:
        p = text_frame.add_paragraph()
        p.text = f"• {point}"
        p.font.size = Pt(15)
        p.font.color.rgb = RGBColor(40, 40, 40)
        p.level = 1
        p.space_after = Pt(3)

# 슬라이드 12: Email 템플릿
slide12 = prs.slides.add_slide(prs.slide_layouts[6])
background = slide12.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(250, 250, 250)

title_box = slide12.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.7))
title_frame = title_box.text_frame
title_frame.text = "Cold Email 템플릿 예시"
title_paragraph = title_frame.paragraphs[0]
title_paragraph.font.size = Pt(32)
title_paragraph.font.bold = True
title_paragraph.font.color.rgb = RGBColor(60, 60, 60)

email_box = slide12.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.3))
text_frame = email_box.text_frame
text_frame.word_wrap = True

email_template = """Subject: Interest in Your Research on [Specific Topic]

Dear Professor [Name],

I am [Your Name], a senior undergraduate student majoring in Computer Science at Kyung Hee University. I recently read your paper "[Paper Title]" published at [Conference/Journal], and I was particularly fascinated by [specific technical aspect].

I have been working on [related project/research area], and I believe my background aligns well with your research. Here are some of my relevant works:

• GitHub: [link] - [Brief description of key project]
• Technical Blog: [link] - [Brief description of deep-dive post]

I would be honored to discuss potential opportunities to join your lab as a graduate student. Would you be available for a brief 15-20 minute meeting?

Thank you for your time and consideration.

Best regards,
[Your Name]
"""

p = text_frame.paragraphs[0]
p.text = email_template
p.font.size = Pt(14)
p.font.name = 'Courier New'
p.font.color.rgb = RGBColor(40, 40, 40)
p.line_spacing = 1.3

# 슬라이드 13: Phase 4 - 석사과정
slide13 = prs.slides.add_slide(prs.slide_layouts[6])
background = slide13.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(243, 229, 245)

title_box = slide13.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
title_frame = title_box.text_frame
title_frame.text = "Phase 4: 석사과정 (2027.03 ~ 2029.02)"
title_paragraph = title_frame.paragraphs[0]
title_paragraph.font.size = Pt(36)
title_paragraph.font.bold = True
title_paragraph.font.color.rgb = RGBColor(106, 27, 154)

subtitle_box = slide13.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(0.6))
subtitle_frame = subtitle_box.text_frame
subtitle_frame.text = "🎯 목표: 연구 성과 창출 + 취업 준비"
subtitle_paragraph = subtitle_frame.paragraphs[0]
subtitle_paragraph.font.size = Pt(24)
subtitle_paragraph.font.bold = True
subtitle_paragraph.font.color.rgb = RGBColor(123, 31, 162)

content_box = slide13.shapes.add_textbox(Inches(0.8), Inches(2.3), Inches(8.4), Inches(4.7))
text_frame = content_box.text_frame
text_frame.word_wrap = True

sections = [
    ("1년차 (2027.03 ~ 2027.12)", [
        "필수 과목 & 고급 전공 수업 수강",
        "석사 연구 주제 선정 및 실험",
        "여름: 인턴십/산학협동 기회 모색",
        "국내 학회 논문 발표 목표"
    ]),
    ("2년차 (2028.01 ~ 2029.02)", [
        "국제 학회 논문 제출 (상반기 1건)",
        "졸업 논문 완성 및 심사",
        "취업 준비 병행 (하반기부터)",
        "오픈소스 기여 & 영어 발표 능력 강화"
    ]),
    ("목표 성과", [
        "국제 학회 발표 1회",
        "국내외 저널/학회 1-2회",
        "석사학위 취득",
        "희망 직장 입사 (2029.03~)"
    ])
]

for i, (section_title, points) in enumerate(sections):
    p = text_frame.add_paragraph() if i > 0 else text_frame.paragraphs[0]
    p.text = f"📌 {section_title}"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RGBColor(106, 27, 154)
    p.space_after = Pt(8)
    if i > 0:
        p.space_before = Pt(14)
    
    for point in points:
        p = text_frame.add_paragraph()
        p.text = f"• {point}"
        p.font.size = Pt(17)
        p.font.color.rgb = RGBColor(40, 40, 40)
        p.level = 1
        p.space_after = Pt(4)

# 슬라이드 14: 추천 프로젝트
slide14 = prs.slides.add_slide(prs.slide_layouts[6])
background = slide14.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(240, 248, 255)

title_box = slide14.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.7))
title_frame = title_box.text_frame
title_frame.text = "추천 프로젝트 (FastAPI + PyTorch + 분산학습)"
title_paragraph = title_frame.paragraphs[0]
title_paragraph.font.size = Pt(30)
title_paragraph.font.bold = True
title_paragraph.font.color.rgb = RGBColor(13, 71, 161)

content_box = slide14.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(8.6), Inches(5.5))
text_frame = content_box.text_frame
text_frame.word_wrap = True

projects = [
    ("1. 대용량 이미지 분류 서비스", "Celery 분산 처리 기반 모델 API\nFastAPI + PyTorch + Docker + Redis"),
    ("2. 분산 딥러닝 훈련 관리 웹앱", "Distributed Training Dashboard\n실시간 학습 모니터링 & 시각화"),
    ("3. Federated Learning 시뮬레이터", "연합학습 클라이언트-서버 시스템\nFlower 프레임워크 + FastAPI"),
    ("4. 실시간 분산 추론 서비스", "Load Balancing을 통한 추론 최적화\n로드밸런싱 + 멀티워커 구조"),
    ("5. 자동 ML 파이프라인 & CI/CD", "MLOps 종합 플랫폼\nAutoML + 자동 배포 + 모니터링")
]

for i, (title, desc) in enumerate(projects):
    p = text_frame.add_paragraph() if i > 0 else text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(13, 71, 161)
    p.space_after = Pt(4)
    if i > 0:
        p.space_before = Pt(14)
    
    for line in desc.split('\n'):
        p = text_frame.add_paragraph()
        p.text = f"  • {line}"
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(60, 60, 60)
        p.level = 1
        p.space_after = Pt(3)

# 슬라이드 15: 6개월 실행 타임라인
slide15 = prs.slides.add_slide(prs.slide_layouts[6])
background = slide15.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(245, 250, 255)

title_box = slide15.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.7))
title_frame = title_box.text_frame
title_frame.text = "YK/IST 대비 6개월 타임라인"
title_paragraph = title_frame.paragraphs[0]
title_paragraph.font.size = Pt(32)
title_paragraph.font.bold = True
title_paragraph.font.color.rgb = RGBColor(1, 87, 155)

content_box = slide15.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(8.4), Inches(5.3))
text_frame = content_box.text_frame
text_frame.word_wrap = True

timeline = [
    ("Month 1-2: 포트폴리오 구축", [
        "GitHub 프로젝트 3개 선정 & README 작성",
        "코드 리팩토링 & 주석 추가",
        "기술 블로그 플랫폼 선정 & 첫 포스트"
    ]),
    ("Month 3-4: 콘텐츠 강화", [
        "기술 블로그 2-3개 추가 포스팅",
        "GitHub 프로젝트 테스트 코드 추가",
        "타겟 대학 연구실 리스트업 (5-10개)"
    ]),
    ("Month 5: 컨택 준비", [
        "교수님 연구 분야 상세 조사",
        "맞춤형 Cold Email 초안 작성",
        "이메일 리뷰 & 피드백"
    ]),
    ("Month 6: 공격적 컨택", [
        "Cold Email 발송 (주 2-3명)",
        "Follow-up 관리",
        "교수님과 미팅 진행",
        "필요 시 Open Lab 지원"
    ])
]

for i, (phase, tasks) in enumerate(timeline):
    p = text_frame.add_paragraph() if i > 0 else text_frame.paragraphs[0]
    p.text = phase
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RGBColor(1, 87, 155)
    p.space_after = Pt(8)
    if i > 0:
        p.space_before = Pt(14)
    
    for task in tasks:
        p = text_frame.add_paragraph()
        p.text = f"✓ {task}"
        p.font.size = Pt(17)
        p.font.color.rgb = RGBColor(40, 40, 40)
        p.level = 1
        p.space_after = Pt(4)

# 슬라이드 16: 핵심 성공 요인
slide16 = prs.slides.add_slide(prs.slide_layouts[6])
background = slide16.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(232, 245, 233)

title_box = slide16.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.7))
title_frame = title_box.text_frame
title_frame.text = "핵심 성공 요인 (Key Success Factors)"
title_paragraph = title_frame.paragraphs[0]
title_paragraph.font.size = Pt(32)
title_paragraph.font.bold = True
title_paragraph.font.color.rgb = RGBColor(27, 94, 32)

content_box = slide16.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(5))
text_frame = content_box.text_frame
text_frame.word_wrap = True

factors = [
    ("Consistency (일관성)", "6개월간 꾸준한 포트폴리오 관리 및 블로그 운영"),
    ("Quality over Quantity", "프로젝트 10개보다 고품질 3개가 더 효과적"),
    ("Specificity (구체성)", "일반적 관심보다 구체적인 연구 주제와 경험"),
    ("Proactiveness (적극성)", "답변 기다리기보다 Follow-up과 대안 모색"),
    ("Continuous Learning", "최신 기술 트렌드 학습 및 적용")
]

for i, (factor, description) in enumerate(factors):
    p = text_frame.add_paragraph() if i > 0 else text_frame.paragraphs[0]
    p.text = f"🎯 {factor}"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(27, 94, 32)
    p.space_after = Pt(6)
    if i > 0:
        p.space_before = Pt(18)
    
    p = text_frame.add_paragraph()
    p.text = description
    p.font.size = Pt(19)
    p.font.color.rgb = RGBColor(40, 40, 40)
    p.level = 1
    p.space_after = Pt(4)

# 슬라이드 17: Next Steps
slide17 = prs.slides.add_slide(prs.slide_layouts[6])
background = slide17.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(20, 30, 48)

title_box = slide17.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
title_frame = title_box.text_frame
title_frame.text = "Next Steps"
title_paragraph = title_frame.paragraphs[0]
title_paragraph.font.size = Pt(48)
title_paragraph.font.bold = True
title_paragraph.font.color.rgb = RGBColor(255, 255, 255)
title_paragraph.alignment = PP_ALIGN.CENTER

action_box = slide17.shapes.add_textbox(Inches(1.5), Inches(3.5), Inches(7), Inches(3))
text_frame = action_box.text_frame
text_frame.word_wrap = True

actions = [
    "1. 부트캠프 기간 집중: 이론 + 실전 프로젝트",
    "2. 학부연구생으로 연구 경험 쌓기",
    "3. 포트폴리오 3개+ 완성 (GitHub + 블로그)",
    "4. 대학원 옵션별 전략 실행",
    "5. 2029년 목표: 석사 학위 + 희망 직장"
]

for i, action in enumerate(actions):
    p = text_frame.add_paragraph() if i > 0 else text_frame.paragraphs[0]
    p.text = action
    p.font.size = Pt(22)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.space_after = Pt(16)

final_box = slide17.shapes.add_textbox(Inches(1), Inches(6.3), Inches(8), Inches(0.8))
text_frame = final_box.text_frame
p = text_frame.paragraphs[0]
p.text = "💪 Consistency + Quality + Proactiveness = Success"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 215, 0)
p.alignment = PP_ALIGN.CENTER

# 저장
output_path = "/home/hyuksu/projects/ml/커리어_로드맵_2025-2029.pptx"
prs.save(output_path)
print(f"✅ 프레젠테이션이 생성되었습니다: {output_path}")
print(f"📊 총 {len(prs.slides)} 슬라이드")
print("\n📌 슬라이드 구성:")
print("  - 전체 로드맵 개요 (2025-2029)")
print("  - Phase 1: 부트캠프 (2025.11~2026.04)")
print("  - Phase 2: 학부연구생 (2026.05~2026.12)")
print("  - Phase 3: 대학원 진학 3가지 옵션")
print("  - Phase 4: 석사과정 (2027~2029)")
print("  - YK/IST 상세 전략 (GitHub, 블로그, Cold Email)")
print("  - 추천 프로젝트 & 실행 타임라인")
